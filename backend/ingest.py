import os
import uuid
import fitz
import docx
import openpyxl
import pptx 
import csv 
import io
import pytesseract
from PIL import Image
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
from chromadb.config import Settings

import config
import database
import ollama_client
import asyncio

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
chroma_client = chromadb.PersistentClient(
    path=config.CHROMA_PERSIST_DIR,
    settings=Settings(anonymized_telemetry=False)
)

collection = chroma_client.get_or_create_collection(
    name="documents",
    metadata={"hnsw:space": "cosine"}
)

splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,
    chunk_overlap=150,
)


def extract_from_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    full_text = []

    for page in doc:
        text = page.get_text().strip()
        if text:
            full_text.append(text)
        else:
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            ocr_text = pytesseract.image_to_string(img).strip()
            if ocr_text:
                full_text.append(ocr_text)

    doc.close()
    return "\n\n".join(full_text)

def extract_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)

def extract_from_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    all_text = []
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_text.append(f"[Sheet: {sheet_name}]")
        
        # Convert to CSV format internally
        output = io.StringIO()
        writer = csv.writer(output)
        
        for row in ws.iter_rows(values_only=True):
            # Only write rows that have at least one non-empty cell
            cleaned_row = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in cleaned_row):
                writer.writerow(cleaned_row)
        
        csv_content = output.getvalue()
        if csv_content.strip():
            all_text.append(csv_content)
    
    return "\n\n".join(all_text)

def extract_from_csv(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_from_pptx(file_path: str) -> str:
    prs = pptx.Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slides_text.append(f"[Slide {i+1}]\n" + "\n".join(slide_content))
    return "\n\n".join(slides_text)


def extract_from_image(file_path: str) -> str:
    img = Image.open(file_path)
    return pytesseract.image_to_string(img).strip()

def extract_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
    
def extract_text(file_path: str, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1] 
    extractors = {
        "pdf" : extract_from_pdf,
        "docx" : extract_from_docx,
        "xlsx" : extract_from_xlsx,
        "pptx" : extract_from_pptx,
        "txt" : extract_from_txt,
        "png" : extract_from_image,
        "jpg" : extract_from_image,
        "jpeg" : extract_from_image,
        "csv" : extract_from_csv,
    }

    if ext not in extractors:
        raise ValueError(f"Unsupported file type: .{ext}")
    return extractors[ext](file_path)


async def ingest_document(
        file_path: str,
        filename: str,
        document_id: str,
        owner_id: int,
        visibility: str,
) -> int:
    """
    Full pipeline: extract -> chunk -> embed -> store.
    Returns the number of chunks indexed.
    Raises on failure so the caller can mark the document as failed
    rather than leaving it in a broken half indexed state.
    """

    text = extract_text(file_path, filename)
    if not text.strip():
        raise ValueError(
            "No text could be extracted from this file. "
            "If it's a scannned document, make sure Tesseract is installed."
        )
    
    chunks = splitter.split_text(text)
    if len(chunks)>50:
        print(f"Document has {len(chunks)} chunks, limiting to 50")
        chunks = chunks[:50]
    if not chunks:
        raise ValueError("Document was too short to produce any chunks.")
    
    ids = []
    embeddings = []
    documents = []
    metadatas = []

    for i, chunk in enumerate(chunks):
        vector = await ollama_client.embed_text(chunk)
        chunk_id = f"{document_id}_chunk_{i}"

        ids.append(chunk_id)
        embeddings.append(vector)
        documents.append(chunk)
        metadatas.append({
            "document_id" : document_id,
            "filename" : filename,
            "owner_id" : owner_id,
            "visibility" : visibility,
            "chunk_index" : i,
        })

    collection.add(
        ids=ids,
        embeddings=embeddings,
        documents=documents,
        metadatas=metadatas,
    )

    return len(chunks)